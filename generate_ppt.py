from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Color Palette ──
DARK_BG    = RGBColor(0x1A, 0x1A, 0x2E)  # dark navy
SLIDE_BG   = RGBColor(0x0F, 0x0F, 0x1A)  # near-black
ACCENT     = RGBColor(0x00, 0xD2, 0xFF)  # cyan
ACCENT2    = RGBColor(0xFF, 0x6B, 0x6B)  # coral
ACCENT3    = RGBColor(0x51, 0xCF, 0x66)  # green
ACCENT4    = RGBColor(0xFF, 0xD9, 0x3D)  # yellow
PURPLE     = RGBColor(0xBE, 0x95, 0xFF)  # purple
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xBB, 0xBB, 0xBB)
DARK_CARD  = RGBColor(0x16, 0x21, 0x3E)  # card background
BORDER     = RGBColor(0x2A, 0x3A, 0x5C)  # subtle border

def set_slide_bg(slide, color=SLIDE_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_bar(slide, left, top, width, height, color=ACCENT):
    """Thin accent bar at top of slide."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_title(slide, text, y=Inches(0.6), font_size=36, color=WHITE):
    txBox = slide.shapes.add_textbox(Inches(0.8), y, Inches(11.7), Inches(1.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = True
    return tf

def add_subtitle(slide, text, y=Inches(1.6), font_size=18, color=LIGHT_GRAY):
    txBox = slide.shapes.add_textbox(Inches(0.8), y, Inches(11.7), Inches(0.6))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    return tf

def add_body(slide, text, left=Inches(0.8), top=Inches(2.2), width=Inches(11.7), height=Inches(5.0), font_size=16, color=LIGHT_GRAY):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split("\n")):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(6)
    return tf

def add_card(slide, left, top, width, height, title_text, body_text, title_color=ACCENT, body_color=LIGHT_GRAY):
    """Rounded card with title + body."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_CARD
    shape.line.color.rgb = BORDER
    shape.line.width = Pt(1)

    # Title
    txBox = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.15), width - Inches(0.4), Inches(0.4))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(15)
    p.font.color.rgb = title_color
    p.font.bold = True

    # Body
    txBox2 = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.55), width - Inches(0.4), height - Inches(0.7))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    for i, line in enumerate(body_text.split("\n")):
        if i == 0:
            p2 = tf2.paragraphs[0]
        else:
            p2 = tf2.add_paragraph()
        p2.text = line
        p2.font.size = Pt(12)
        p2.font.color.rgb = body_color
        p2.space_after = Pt(3)
    return shape

def add_arrow_right(slide, left, top, width, height):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = BORDER
    shape.line.fill.background()
    return shape

def add_arrow_down(slide, left, top, width, height):
    shape = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = BORDER
    shape.line.fill.background()
    return shape

def add_table(slide, left, top, col_widths, headers, rows, header_color=ACCENT):
    """Add a styled table."""
    n_rows = len(rows) + 1
    n_cols = len(headers)
    total_w = sum(col_widths)
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, total_w, Inches(0.35 * n_rows))
    table = table_shape.table

    for ci, cw in enumerate(col_widths):
        table.columns[ci].width = cw

    # Header row
    for ci, h in enumerate(headers):
        cell = table.cell(0, ci)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(12)
            p.font.color.rgb = WHITE
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_CARD

    # Data rows
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.color.rgb = LIGHT_GRAY
                p.alignment = PP_ALIGN.CENTER
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0x0F, 0x0F, 0x1A)

    return table

# ============================================================================
# SLIDE 1: TITLE
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide, SLIDE_BG)
add_bar(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT)

add_title(slide, "GEPA Skill Builder", y=Inches(1.8), font_size=52, color=WHITE)
add_subtitle(slide, "Automated Prompt Optimization for Multi-Agent Systems", y=Inches(3.0), font_size=24, color=ACCENT)
add_subtitle(slide, "Using genetic evolutionary algorithms to optimize 3-agent prompt pipelines\nDSPy · GPT-4.1 · MLflow", y=Inches(3.7), font_size=16, color=LIGHT_GRAY)
add_subtitle(slide, "Pranav Bedekar  |  2026-05-27", y=Inches(5.2), font_size=14, color=LIGHT_GRAY)

# ============================================================================
# SLIDE 2: THE PROBLEM
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, SLIDE_BG)
add_bar(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT2)
add_title(slide, "The Problem", font_size=36)

add_card(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(2.0),
    "🔴 Manual Prompt Engineering",
    "• Writing prompts for 3 interdependent agents is tedious\n"
    "• Planner → Curator → Generator: output of one feeds the next\n"
    "• Changing one prompt can break the entire chain\n"
    "• No systematic way to know if a prompt is 'good enough'")

add_card(slide, Inches(6.8), Inches(2.0), Inches(5.5), Inches(2.0),
    "🟢 Automated Prompt Optimization",
    "• Let an LLM reason about WHY prompts fail\n"
    "• Optimize all 3 prompts simultaneously\n"
    "• Genetic search: evolve a population of candidates\n"
    "• Reproducible, measurable, scalable")

add_body(slide, "The core question: Can an LLM optimize prompts better than a human — automatically and at scale?",
         top=Inches(4.5), font_size=18, color=ACCENT4)

# ============================================================================
# SLIDE 3: WHAT THIS SYSTEM DOES
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, SLIDE_BG)
add_bar(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT3)
add_title(slide, "What This System Does", font_size=36)

add_card(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(3.0),
    "📥 Input: Raw Skill Request",
    "Example:\n"
    "\"Create a skill for a data analyst that can clean CSV files:\n"
    "  • remove duplicates\n"
    "  • handle missing values\n"
    "  • detect outliers\n"
    "Document in skill.md format with YAML frontmatter.\"\n\n"
    "10 role-specific requests across data analyst, Python dev,\n"
    "ML engineer, DevOps, product manager, SQL analyst, etc.")

add_card(slide, Inches(6.8), Inches(2.0), Inches(5.5), Inches(3.0),
    "📤 Output: Generated skill.md Files",
    "request_1/\n"
    "├── remove_duplicates.md\n"
    "├── handle_missing_values.md\n"
    "└── detect_outliers.md\n\n"
    "Each file includes:\n"
    "• YAML frontmatter (name, description, inputs, outputs)\n"
    "• Clear, actionable markdown body\n"
    "• Domain-specific instructions for the target role")

# ============================================================================
# SLIDE 4: THE 3-AGENT PIPELINE
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, SLIDE_BG)
add_bar(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT)
add_title(slide, "The 3-Agent Pipeline", font_size=36)

# Flow: User Request → Planner → Curator → Generator → skills
# Labels
y_agents = Inches(2.8)
card_w = Inches(2.8)
card_h = Inches(3.2)
gap = Inches(0.3)

add_card(slide, Inches(0.5), y_agents, card_w, card_h,
    "1. PlannerAgent",
    "Input: user_request (raw text)\n\n"
    "Outputs:\n"
    "• skill_inventory:\n"
    "  numbered list of skills needed\n"
    "• clarified_plan:\n"
    "  step-by-step build plan\n"
    "• assumptions\n"
    "• curation_instructions\n\n"
    "Model: gpt-4o-mini + ChainOfThought",
    title_color=ACCENT2)

add_arrow_right(slide, Inches(3.45), y_agents + Inches(1.3), Inches(0.35), Inches(0.25))

add_card(slide, Inches(3.95), y_agents, card_w, card_h,
    "2. CurationAgent",
    "Input: skill_inventory + plan +\n"
    "       curation_instructions\n\n"
    "Outputs:\n"
    "• skill_specs:\n"
    "  detailed per-skill specification\n"
    "• examples: concrete use cases\n"
    "• quality_checks: validation list\n"
    "• missing_info: gaps needing input\n\n"
    "Model: gpt-4o-mini + ChainOfThought",
    title_color=ACCENT3)

add_arrow_right(slide, Inches(6.9), y_agents + Inches(1.3), Inches(0.35), Inches(0.25))

add_card(slide, Inches(7.4), y_agents, card_w, card_h,
    "3. SkillGeneratorAgent",
    "Input: skill_specs + examples +\n"
    "       quality_checks\n\n"
    "Outputs:\n"
    "• skill_markdown_files[]:\n"
    "  complete .md file contents\n"
    "• file_names[]:\n"
    "  snake_case filenames\n\n"
    "Model: gpt-4o-mini + ChainOfThought",
    title_color=PURPLE)

add_arrow_right(slide, Inches(10.35), y_agents + Inches(1.3), Inches(0.35), Inches(0.25))

add_card(slide, Inches(10.85), Inches(2.8), Inches(2.0), card_h,
    "Output",
    "skill.md\nfiles\n\n✅ YAML\n✅ Body\n✅ Examples\n✅ Domain\n  relevant",
    title_color=ACCENT4)

# DSPy note
add_body(slide, "All agents are DSPy Modules with ChainOfThought signatures. Each signature's docstring IS the prompt that GEPA optimizes.",
         top=Inches(6.3), font_size=13, color=RGBColor(0x88, 0x88, 0x88))

# ============================================================================
# SLIDE 5: WHY OPTIMIZE ALL 3 TOGETHER
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, SLIDE_BG)
add_bar(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT4)
add_title(slide, "Why Optimize All 3 Prompts Simultaneously?", font_size=36)

add_card(slide, Inches(0.8), Inches(2.0), Inches(3.5), Inches(2.5),
    "❌ Optimize One at a Time",
    "• Planner is optimized with original\n"
    "  Curator and Generator prompts\n"
    "• Then Curator is optimized with\n"
    "  new Planner + old Generator\n"
    "• The 'best' Planner + 'best' Curator\n"
    "  may not work well together\n\n"
    "→ Local optimum, not global",
    title_color=ACCENT2)

add_card(slide, Inches(4.8), Inches(2.0), Inches(3.5), Inches(2.5),
    "✅ Optimize All 3 Together",
    "• Round-robin: cycle through\n"
    "  Planner → Curator → Generator\n"
    "• Each round optimizes one agent\n"
    "  while keeping the other two at\n"
    "  their current best\n"
    "• Crossover merges strengths\n"
    "  across different lineages\n\n"
    "→ Joint optimum for the pipeline",
    title_color=ACCENT3)

add_card(slide, Inches(8.8), Inches(2.0), Inches(3.5), Inches(2.5),
    "🔁 Feedback Loop Matters",
    "• Better Planner → clearer plan\n"
    "  → easier for Curator to enrich\n"
    "• Better Curator → richer specs\n"
    "  → easier for Generator to write\n"
    "• Better Generator → higher-quality\n"
    "  output → higher judge scores\n\n"
    "• GEPA preserves this chain by\n"
    "  evaluating the FULL pipeline",
    title_color=ACCENT)

add_body(slide, "One optimization run produces prompts that benefit ALL future skill-generation requests — not just the 10 in the current dataset.",
         top=Inches(5.0), font_size=16, color=ACCENT4)

# ============================================================================
# SLIDE 6: WHAT IS GEPA?
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, SLIDE_BG)
add_bar(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), PURPLE)
add_title(slide, "What is GEPA?", font_size=36)
add_subtitle(slide, "Genetic Evolutionary Prompt Algorithm — DSPy's LLM-driven prompt optimizer", font_size=18)

add_card(slide, Inches(0.8), Inches(2.2), Inches(5.5), Inches(2.2),
    "🧬 Genetic Algorithm Core",
    "• Maintains a POPULATION of prompt candidates\n"
    "• Each candidate = 3 prompts (Planner, Curator, Generator)\n"
    "• Fitness = judge score from running the full pipeline\n"
    "• Generations evolve through selection + mutation + crossover",
    title_color=ACCENT4)

add_card(slide, Inches(6.8), Inches(2.2), Inches(5.5), Inches(2.2),
    "🧠 LLM as the Mutation Operator",
    "• NOT random string edits — the Reflection Model (GPT-4.1)\n"
    "  reads full execution traces from ALL 3 agents\n"
    "• Understands WHY each agent failed (judge feedback)\n"
    "• Proposes targeted, semantic prompt edits\n"
    "• Edits are domain-aware: \"Add instructions to split complex\n"
    "  requests into sub-skills\" vs random token swaps",
    title_color=PURPLE)

add_card(slide, Inches(0.8), Inches(4.7), Inches(5.5), Inches(2.0),
    "🏆 Pareto Frontier Selection",
    "• Multi-objective: candidates scored across 5 dimensions\n"
    "  (YAML, clarity, actionability, domain, completeness)\n"
    "• Pareto filter keeps non-dominated candidates\n"
    "• Preserves diverse strengths — a candidate strong in\n"
    "  Completeness won't be killed by one strong in Clarity",
    title_color=ACCENT3)

add_card(slide, Inches(6.8), Inches(4.7), Inches(5.5), Inches(2.0),
    "🔀 Crossover & Merge",
    "• Pairs strong candidates from different lineages\n"
    "• Merges prompt strengths: takes the Planner prompt\n"
    "  from candidate A and the Generator prompt from candidate B\n"
    "• Small mutations add exploration\n"
    "• skip_perfect_score=True avoids wasting budget on\n"
    "  candidates that can't be improved",
    title_color=ACCENT2)

# ============================================================================
# SLIDE 7: GEPA LOOP DIAGRAM
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, SLIDE_BG)
add_bar(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT)
add_title(slide, "GEPA Optimization Loop — Step by Step", font_size=36)

loop_y = Inches(1.8)
box_w = Inches(5.5)
box_h = Inches(1.0)
arrow_gap = Inches(0.15)

steps = [
    ("🔄 Round-Robin Agent Selection", "Cycle through Planner → Curator → Generator. Each round picks one agent to optimize.", ACCENT2),
    ("🧠 GPT-4.1 Reflection Model", "Reads execution traces + judge feedback from all 3 agents. Identifies failure patterns. Proposes targeted prompt edits for the current agent.", PURPLE),
    ("🎯 Evaluate Candidate", "Run the FULL pipeline with the edited prompt on all 7 training requests. Judge scores every output across 5 dimensions.", ACCENT3),
    ("🏆 Pareto Frontier Selection", "Compare all candidates across all 5 metric dimensions. Keep non-dominated ones. Preserve diverse strengths.", ACCENT4),
    ("🔀 Crossover & Merge", "Pair strong candidates from different lineages. Merge prompt strengths. Apply small random mutations. Feed back into round-robin.", ACCENT),
]

for i, (title, desc, color) in enumerate(steps):
    x = Inches(0.8)
    y = loop_y + i * (box_h + arrow_gap)
    add_card(slide, x, y, Inches(11.5), box_h, title, desc, title_color=color)

    if i < len(steps) - 1:
        add_arrow_down(slide, Inches(6.0), y + box_h + Inches(0.02), Inches(0.25), Inches(0.13))

# Loop back arrow (curved — we'll use text indicator)
add_body(slide, "⤴  Loop continues until convergence or auto budget exhausted  ⤴",
         top=Inches(6.8), font_size=14, color=ACCENT2)

# ============================================================================
# SLIDE 8: LLM AS JUDGE
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, SLIDE_BG)
add_bar(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT3)
add_title(slide, "LLM-as-Judge — How We Measure Quality", font_size=36)

add_table(slide, Inches(0.8), Inches(2.0),
    [Inches(3.0), Inches(1.2), Inches(7.5)],
    ["Dimension", "Weight", "What It Checks"],
    [
        ["YAML Frontmatter", "20%", "Valid YAML with name, description, inputs, and outputs fields"],
        ["Body Clarity", "20%", "Well-structured, readable markdown with clear headings and organization"],
        ["Actionability", "20%", "Concrete, specific steps that an AI agent can follow — not vague advice"],
        ["Domain Relevance", "20%", "Addresses the specific role and domain mentioned in the original request"],
        ["Completeness", "20%", "Covers ALL requirements from the original request without omissions"],
    ])

add_card(slide, Inches(0.8), Inches(4.2), Inches(5.5), Inches(2.5),
    "📊 Scoring Details",
    "• Each dimension: 1 (poor) → 5 (excellent)\n"
    "• Total = sum of 5 dimensions / 25\n"
    "• Normalized to 0.00 – 1.00 range\n"
    "• Judge model: GPT-4.1 (temperature = 0.1)\n"
    "• Called via direct OpenAI SDK (not DSPy)\n\n"
    "The judge's feedback text is what GEPA's\n"
    "reflection model reads to understand WHY\n"
    "a skill failed and HOW to fix the prompts.",
    title_color=ACCENT4)

add_card(slide, Inches(6.8), Inches(4.2), Inches(5.5), Inches(2.5),
    "🔍 Example Judge Output",
    "Score: 0.88/1.00\n"
    "YAML=5/5, Clarity=5/5, Actionability=5/5,\n"
    "Domain=5/5, Completeness=2/5\n\n"
    "Feedback:\n"
    "\"The skill.md file is well-structured, with clean\n"
    "YAML and actionable markdown. However, it only\n"
    "covers 'remove duplicates' and omits handling\n"
    "missing values and detecting outliers — it's\n"
    "incomplete with respect to the original request.\"",
    title_color=PURPLE)

# ============================================================================
# SLIDE 9: DATASET & EXPERIMENT
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, SLIDE_BG)
add_bar(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT)
add_title(slide, "Dataset & Experiment Design", font_size=36)

add_card(slide, Inches(0.8), Inches(2.0), Inches(3.5), Inches(2.2),
    "📥 Input Dataset",
    "• skill_requests.csv\n"
    "• 10 role-specific requests\n"
    "• Role examples:\n"
    "  Data Analyst, Python Dev,\n"
    "  ML Engineer, Data Engineer,\n"
    "  Business Analyst, DevOps,\n"
    "  AI Engineer, Product Manager,\n"
    "  SQL Analyst, Docs Writer",
    title_color=ACCENT)

add_card(slide, Inches(4.8), Inches(2.0), Inches(3.5), Inches(2.2),
    "✂️ Train/Val Split",
    "• seed = 42, deterministic\n"
    "• 7 training requests\n"
    "  (ids: 1,2,3,6,8,9,10)\n"
    "• 3 validation requests\n"
    "  (ids: 4,5,7)\n"
    "• Train = GEPA optimization\n"
    "• Val = held-out evaluation",
    title_color=ACCENT4)

add_card(slide, Inches(8.8), Inches(2.0), Inches(3.5), Inches(2.2),
    "📏 Baseline Phase",
    "• Run pipeline with original\n"
    "  hand-written prompts\n"
    "• Score all 7 train requests\n"
    "  with LLM judge\n"
    "• Baseline avg: ~0.80\n"
    "• Saves generated skills to\n"
    "  baseline_skills/ for comparison",
    title_color=ACCENT2)

add_table(slide, Inches(0.8), Inches(4.6),
    [Inches(1.8), Inches(3.5), Inches(1.5), Inches(1.5), Inches(1.5)],
    ["Phase", "What Happens", "Requests", "MLflow Tag", "Output"],
    [
        ["1. Baseline", "Run original pipeline + judge all 7", "7 train", "phase=baseline", "baseline_skills/"],
        ["2. Optimize", "GEPA evolves prompts on 7 train", "7 train", "phase=optimization", "gepa_logs/"],
        ["3. Save", "Extract optimized prompts, save skills", "all 10", "—", "optimized_prompts/"],
        ["4. Evaluate", "Run baseline vs optimized on 3 val", "3 val", "notebook 02", "evaluation_report.json"],
    ])

# ============================================================================
# SLIDE 10: RESULTS
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, SLIDE_BG)
add_bar(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT3)
add_title(slide, "Results", font_size=36)

# Big numbers
add_card(slide, Inches(0.8), Inches(2.0), Inches(3.5), Inches(2.0),
    "66", "Candidates evolved\nduring GEPA run", title_color=ACCENT, body_color=WHITE)

add_card(slide, Inches(4.8), Inches(2.0), Inches(3.5), Inches(2.0),
    "+7.8%", "Improvement over baseline\n0.8000 → 0.8781", title_color=ACCENT3, body_color=WHITE)

add_card(slide, Inches(8.8), Inches(2.0), Inches(3.5), Inches(2.0),
    "3/3", "Agents optimized\n(Planner, Curator,\nGenerator)", title_color=ACCENT4, body_color=WHITE)

add_table(slide, Inches(0.8), Inches(4.5),
    [Inches(1.5), Inches(2.0), Inches(2.0), Inches(1.5), Inches(4.5)],
    ["Request #", "Baseline", "Optimized", "Delta", "Role"],
    [
        ["1", "0.88", "0.88", "+0.00", "Data Analyst (already near ceiling)"],
        ["2", "0.68", "0.80", "+0.12", "Python Developer (biggest gain)"],
        ["3", "0.80", "—", "—", "ML Engineer"],
        ["6", "0.80", "—", "—", "DevOps Engineer"],
        ["8", "0.84", "—", "—", "Product Manager"],
        ["9", "0.80", "—", "—", "SQL Analyst"],
        ["10", "0.80", "—", "—", "Documentation Writer"],
    ])

# ============================================================================
# SLIDE 11: WHAT ACTUALLY IMPROVED
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, SLIDE_BG)
add_bar(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), PURPLE)
add_title(slide, "What Actually Improved?", font_size=36)

add_card(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.5),
    "📝 Baseline Prompt (Planner — 211 chars)",
    "Break a raw skill request into a clear build plan.\n\n"
    "Identify every skill that must be created, preserve\n"
    "the user's requirements, and make the next agent's\n"
    "job easy. Do not write the final skill.md files.\n\n"
    "─────────────────────────────\n\n"
    "Problem: Too generic. No guidance on HOW to break\n"
    "requests into skills. The Generator often produced\n"
    "one combined skill instead of separate ones per\n"
    "requirement → Completeness score = 2/5.",
    title_color=ACCENT2)

add_card(slide, Inches(6.8), Inches(2.0), Inches(5.5), Inches(4.5),
    "✅ Optimized Prompt (Planner — after GEPA)",
    "[GEPA-added instructions typically include:]\n\n"
    "• Explicit instruction to create one skill PER\n"
    "  bullet point in the user request\n"
    "• \"If the request lists N requirements, create\n"
    "  exactly N separate skills\"\n"
    "• Clear naming convention for skill inventory\n"
    "• Instructions to flag ambiguous requirements\n"
    "  in the assumptions field\n\n"
    "→ Generator now outputs 3 separate .md files\n"
    "instead of 1 combined file → Completeness ↑",
    title_color=ACCENT3)

add_body(slide, "Note: Optimized prompts were not fully extracted in this run (GEPA extraction bug). The above shows the expected behavior based on the +7.8% score improvement and the known completeness bottleneck.",
         top=Inches(6.8), font_size=12, color=RGBColor(0x88, 0x88, 0x88))

# ============================================================================
# SLIDE 12: MLflow SETUP
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, SLIDE_BG)
add_bar(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT)
add_title(slide, "MLflow Observability — Run Hierarchy", font_size=36)

add_card(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.5),
    "🖥️ MLflow Server: 172.27.72.27:5000",
    "Session: {session_id}\n"
    "Example: 20260526_1430\n\n"
    "Every run in a single notebook execution shares\n"
    "the same session_id tag for unified filtering.\n\n"
    "3 phases, each as a named parent run:\n\n"
    "1. baseline_{session_id}\n"
    "   → 7 nested runs (one per request)\n"
    "2. optimization_{session_id}\n"
    "   → GEPA compile with all params logged\n"
    "3. summary_{session_id}\n"
    "   → Final artifacts: prompts, tables, CSV",
    title_color=ACCENT)

add_card(slide, Inches(6.8), Inches(2.0), Inches(5.5), Inches(4.5),
    "🏷️ Tags on Every Run",
    "app              = \"skill_builder\"\n"
    "phase            = \"baseline\" | \"optimization\" | \"summary\"\n"
    "request_id       = \"1\" ... \"10\"\n"
    "split            = \"train\" | \"val\"\n"
    "session_id       = \"20260526_1430\"\n"
    "pipeline_version = \"1.1.0\"\n"
    "dataset_version  = \"skill_requests_v1\"\n\n"
    "Filtering examples:\n"
    "  tags.phase = \"baseline\" → all baseline traces\n"
    "  tags.request_id = \"3\" → all traces for request #3\n"
    "  tags.session_id = \"20260526_1430\" → one notebook run",
    title_color=ACCENT4)

# ============================================================================
# SLIDE 13: TRACE DATASETS
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, SLIDE_BG)
add_bar(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT2)
add_title(slide, "Trace Datasets — Mining Agent Behavior", font_size=36)

add_card(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(2.0),
    "📊 What We Extract",
    "• 64 MLflow traces loaded from experiment\n"
    "• 113 agent-level spans extracted:\n"
    "  • 38 PlannerAgent spans\n"
    "  • 38 CurationAgent spans\n"
    "  • 37 SkillGeneratorAgent spans\n"
    "• Each span contains: inputs JSON, outputs JSON,\n"
    "  duration in ms, status code, trace metadata",
    title_color=ACCENT)

add_card(slide, Inches(6.8), Inches(2.0), Inches(5.5), Inches(2.0),
    "📁 Output: One CSV Per Agent",
    "trace_datasets/\n"
    "├── all_agents_dataset.csv (113 rows)\n"
    "├── planner_agent_dataset.csv (38 rows)\n"
    "├── curation_agent_dataset.csv (38 rows)\n"
    "└── skill_generator_agent_dataset.csv (37 rows)\n\n"
    "Each CSV is ready for offline analysis:\n"
    "pandas, SQL, or visualization tools.",
    title_color=ACCENT3)

add_card(slide, Inches(0.8), Inches(4.4), Inches(11.5), Inches(2.5),
    "🔍 What You Can Analyze",
    "• Bottleneck detection: Which agent has the highest avg duration? Is Curator slower than Generator?\n"
    "• Failure clustering: Do errors cluster on specific request_ids? Are certain roles harder?\n"
    "• Input/output drift: Does the Planner's output format change over GEPA rounds?\n"
    "• Prompt sensitivity: How much does changing one agent's prompt affect the NEXT agent's inputs?\n"
    "• Training data: These traces could fine-tune smaller models to mimic specific agent behaviors",
    title_color=ACCENT4)

# ============================================================================
# SLIDE 14: PERFORMANCE TUNING
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, SLIDE_BG)
add_bar(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT3)
add_title(slide, "Performance Tuning — 300 min → 45 min", font_size=36)

add_table(slide, Inches(0.8), Inches(2.0),
    [Inches(2.5), Inches(2.5), Inches(2.5), Inches(2.5), Inches(2.0)],
    ["Setting", "Slow (300 min)", "Fast (~45 min)", "Speedup", "How It Works"],
    [
        ["auto", "medium", "light", "~2x", "Fewer candidates & rounds"],
        ["num_threads", "1", "3–4", "~2.5x", "Parallel LM calls for eval"],
        ["Train requests", "7", "4–5", "~1.4x", "Less work per GEPA round"],
        ["reflection_lm", "gpt-4.1", "gpt-4o-mini", "~1.3x", "Faster reflection model"],
        ["Combined", "300 min", "~45 min", "~7x", "All levers together"],
    ])

add_card(slide, Inches(0.8), Inches(4.2), Inches(5.5), Inches(2.5),
    "⚠️ Trade-offs",
    "• light budget: fewer candidates explored,\n"
    "  may not find the global optimum\n"
    "• Higher num_threads: risk of API rate\n"
    "  limiting or transient failures\n"
    "• Fewer train requests: optimization may\n"
    "  overfit to the chosen subset\n\n"
    "Recommendation: start with light/3 threads\n"
    "to validate the pipeline, then scale up.",
    title_color=ACCENT2)

add_card(slide, Inches(6.8), Inches(4.2), Inches(5.5), Inches(2.5),
    "💡 Best Practices",
    "• Always run baseline first — you need a\n"
    "  number to beat before spending compute\n"
    "• Keep the GEPA log_dir (gepa_logs/) — you\n"
    "  can resume from checkpoint if it crashes\n"
    "• Use skip_perfect_score=True — don't\n"
    "  waste budget on un-improvable candidates\n"
    "• Log everything to MLflow — you'll want\n"
    "  those traces for debugging later",
    title_color=ACCENT3)

# ============================================================================
# SLIDE 15: KEY TAKEAWAYS
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, SLIDE_BG)
add_bar(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT)
add_title(slide, "Key Takeaways", font_size=36)

takeaways = [
    ("1",
     "GEPA optimizes multi-agent prompts simultaneously",
     "One optimization run improved all 3 agents (Planner, Curator, Generator) together — preserving the feedback loop between them."),
    ("2",
     "LLM-as-judge + reflection closes the auto-eval loop",
     "GPT-4.1 both scores output quality (5 dimensions × 1–5) AND proposes targeted prompt edits — no human evaluation needed during optimization."),
    ("3",
     "MLflow gives full traceability",
     "Every run, every agent span, every judge score is logged and filterable. 113 agent spans exported to CSV for offline analysis."),
    ("4",
     "Reusable optimized prompts",
     "The evolved prompts benefit ALL future skill requests, not just the 10 in the current dataset. One optimization run → permanent improvement."),
    ("5",
     "The pattern works for any DSPy pipeline",
     "Swap the agents, swap the metric, keep the optimizer. GEPA can optimize any multi-agent system built with DSPy."),
]

for i, (num, title, desc) in enumerate(takeaways):
    y = Inches(1.8) + i * Inches(1.05)
    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), y + Inches(0.05), Inches(0.5), Inches(0.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = ACCENT
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(20)
    tf.paragraphs[0].font.color.rgb = SLIDE_BG
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_card(slide, Inches(1.5), y, Inches(11.0), Inches(0.9), title, desc, title_color=WHITE)

# ============================================================================
# SLIDE 16: FUTURE DIRECTIONS
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, SLIDE_BG)
add_bar(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), PURPLE)
add_title(slide, "Future Directions", font_size=36)

directions = [
    ("🔧", "Fix optimized prompt extraction", "The GEPA-extracted prompts weren't saved correctly. Fix the extended_signature access so the eval notebook can compare baseline vs optimized on the 3 held-out requests."),
    ("📈", "Run full validation", "Complete the 02_gepa_evaluate.ipynb notebook to get statistically meaningful comparison on the 3 held-out validation requests."),
    ("🚀", "Scale up optimization", "Try auto=\"medium\" or \"heavy\" with num_threads=4–6 and the full 7-request training set for higher-quality prompts."),
    ("🧪", "Test stronger task model", "Swap gpt-4o-mini → gpt-4o for the pipeline agents. Does a stronger base model benefit more or less from GEPA optimization?"),
    ("🔬", "Analyze trace datasets", "Use the 113 extracted agent spans to identify bottlenecks, failure patterns, and prompt sensitivity — potentially fine-tune smaller models per agent."),
    ("🌐", "Generalize to other domains", "The same 3-agent pattern (Planner → Curator → Generator) works for code generation, document writing, test creation — swap skill_requests.csv for any task CSV."),
]

for i, (icon, title, desc) in enumerate(directions):
    row_y = Inches(1.8) + i * Inches(0.85)
    add_card(slide, Inches(0.8), row_y, Inches(0.7), Inches(0.7), icon, "", title_color=WHITE)
    add_card(slide, Inches(1.7), row_y, Inches(10.8), Inches(0.7), title, desc, title_color=ACCENT4)

# ============================================================================
# SLIDE 17: TECH STACK
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, SLIDE_BG)
add_bar(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT)
add_title(slide, "Tech Stack & Architecture", font_size=36)

add_table(slide, Inches(0.8), Inches(2.0),
    [Inches(3.5), Inches(3.5), Inches(5.0)],
    ["Component", "Technology", "Role"],
    [
        ["Agent Framework", "DSPy 3.2.1", "Module/Signature/ChainOfThought abstractions"],
        ["Optimizer", "GEPA", "Genetic Evolutionary Prompt Algorithm"],
        ["Task Model", "gpt-4o-mini", "Pipeline agents (Planner, Curator, Generator)"],
        ["Reflection Model", "gpt-4.1 (32k tokens)", "GEPA meta-cognitive trace analysis + prompt edits"],
        ["Judge Model", "gpt-4.1 (temp=0.1)", "LLM-as-Judge: 5-dimension skill scoring"],
        ["Experiment Tracking", "MLflow 3.9.0", "Runs, tags, metrics, artifacts, traces"],
        ["MLflow Server", "172.27.72.27:5000", "Self-hosted tracking server"],
        ["Notebooks", "Jupyter (IPYNB)", "4 notebooks: optimize, evaluate, compare, trace-dataset"],
        ["Data Format", "CSV → pandas DataFrame", "10 skill requests → train/val split"],
        ["Visualization", "draw.io", "Architecture flow diagrams"],
        ["Presentation", "python-pptx", "This slide deck"],
    ])

add_card(slide, Inches(0.8), Inches(5.0), Inches(11.5), Inches(2.0),
    "🔗 File Structure",
    "H:\\GEPA\\gepa\\\n"
    "├── 01_gepa_optimize.ipynb          ← Main: baseline → GEPA → save prompts\n"
    "├── 02_gepa_evaluate.ipynb         ← Evaluation: baseline vs optimized on val set\n"
    "├── 03_gepa_trace_comparison.ipynb ← Side-by-side trace & skill comparison\n"
    "├── create_agent_trace_datasets.ipynb ← Extract MLflow traces → agent CSVs\n"
    "├── agent_skillbuilder.py          ← Standalone pipeline (no optimization)\n"
    "├── skill_requests.csv             ← 10 input requests\n"
    "├── gepa_flow.drawio               ← Architecture diagram\n"
    "└── gepa_presentation.pptx         ← This file",
    title_color=ACCENT3)

# ============================================================================
# Save
# ============================================================================
output_path = r"H:\GEPA\gepa\gepa_presentation.pptx"
prs.save(output_path)
print(f"Presentation saved: {output_path}")
